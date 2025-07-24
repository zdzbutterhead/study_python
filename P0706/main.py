'''
# 输入m和n，计算组合数C(m,n)的值
m = int(input("Enter the value of m: "))
n = int(input("Enter the value of n: "))

fm = 1
fn = 1
fn_m = 1

for i in range(1, m+1):
    fm*=i
for i in range(1, n+1):
    fn*=i
for i in range(1,m-n+1):
    fn_m*=i

C = fm//fn//fn_m
print("The value of C(m,n) is:", C)
'''


'''
# 输入m和n，计算组合数C(m,n)的值
m = int(input("Enter the value of m: "))
n = int(input("Enter the value of n: "))

def factorial(n):
    if n == 0 or n == 1:
        return 1
    else:
        return n * factorial(n-1)

fm = factorial(m)
fn = factorial(n)
fn_m = factorial(m-n)

print("The value of C(m,n) is:",fm//fn//fn_m)
'''


'''
# 设计一个生成随机验证码的函数，验证码由数字和英文大小写字母构成，长度可以通过参数设置。
import random

def generate_code(n = 4):
    # 小写字母
    lalp = [chr(i) for i in range(ord('a'), ord('z')+1)]
    # 大写字母
    salp = [chr(i) for i in range(ord('A'), ord('Z')+1)]
    # 合并为所有字母
    alp = lalp + salp
    # 数字
    num = [str(i) for i in range(10)]
    
    # 随机生成n个字母
    anum = random.randint(0,n)
    # 剩余的数字个数
    nnum = n - anum

    # 随机生成验证码
    code = random.sample(alp,anum)
    code += random.sample(num,nnum)

    # 利用集合的无序性来打乱验证码
    set_code = set(code)
    # 输出打乱后的验证码
    for i in set_code:
        print(i, end='')


for i in range(5):
    generate_code(5)
    print()
'''


'''
def prime(n):
    if n <2:
        return False
    for i in range(2,int(n**0.5)+1):
        if n % i == 0:
            return False
    return True
'''


'''
def greatest_common_divisor(a: int, b: int):
    if a < b:
        while b % a != 0:
            a, b = b % a, a
        return a
    else:
        while a % b != 0:
            a, b = b, a % b
        return b

def least_common_multiple(a: int, b: int):
    return a * b // greatest_common_divisor(a, b)

a=greatest_common_divisor(48, 18)
print(a)

b=least_common_multiple(48, 18)
print(b)
'''


'''
def sample_mean(data):
    return sum(data)/len(data)

def median(data):
    data.sort()
    if len(data) % 2 == 0:
        return (data[len(data)//2] + data[len(data)//2 - 1])/2
    else:
        return data[len(data)//2]

def ptp(data):
    return max(data) - min(data)

def sample_variance(data):
    return sum((x - sample_mean(data))**2 for x in data)/(len(data)-1)

def sample_standard_deviation(data):
    return pow(sample_variance(data),0.5)

def coefficient_of_sample_variation(data):
    return sample_standard_deviation(data)/sample_mean(data)
'''


'''
old_strings = ['in', 'apple', 'zoo', 'waxberry', 'pear']
old_strings.sort(key=len)
print(old_strings)  # ['in', 'zoo', 'pear', 'apple', 'waxberry']
'''


'''
from enum import Enum


class Suite(Enum):
    """花色(枚举)"""
    SPADE, HEART, CLUB, DIAMOND = range(4)


class Card:
    """牌"""

    def __init__(self, suite, face):
        self.suite = suite
        self.face = face

    def __repr__(self):
        suites = '♠♥♣♦'
        faces = ['', 'A', '2', '3', '4', '5', '6', '7', '8', '9', '10', 'J', 'Q', 'K']
        return f'{suites[self.suite.value]}{faces[self.face]}'

    def __lt__(self, other):
        if self.suite == other.suite:
            return self.face < other.face  # 花色相同比较点数的大小
        return self.suite.value < other.suite.value  # 花色不同比较花色对应的值

import random


class Poker:
    """扑克"""

    def __init__(self):
        self.cards = [Card(suite, face)
                      for suite in Suite
                      for face in range(1, 14)]  # 52张牌构成的列表
        self.current = 0  # 记录发牌位置的属性

    def shuffle(self):
        """洗牌"""
        self.current = 0
        random.shuffle(self.cards)  # 通过random模块的shuffle函数实现随机乱序

    def deal(self):
        """发牌"""
        card = self.cards[self.current]
        self.current += 1
        return card

    @property
    def has_next(self):
        """还有没有牌可以发"""
        return self.current < len(self.cards)

class Player:
    """玩家"""

    def __init__(self, name):
        self.name = name
        self.cards = []  # 玩家手上的牌

    def get_one(self, card):
        """摸牌"""
        self.cards.append(card)

    def arrange(self):
        """整理手上的牌"""
        self.cards.sort()

poker = Poker()
poker.shuffle()
players = [Player('东邪'), Player('西毒'), Player('南帝'), Player('北丐')]
# 将牌轮流发到每个玩家手上每人13张牌
for _ in range(13):
    for player in players:
        player.get_one(poker.deal())
# 玩家整理手上的牌输出名字和手牌
for player in players:
    player.arrange()
    print(f'{player.name}: ', end='')
    print(player.cards)
'''


'''
file = open('致橡树.txt', 'r', encoding='utf-8')
for line in file:
    print(line, end='')
file.close()

file = open('致橡树.txt', 'r', encoding='utf-8')
lines = file.readlines()
for line in lines:
    print(line, end='')
file.close()
'''


'''
file = open('致橡树.txt', 'a', encoding='utf-8')
file.write('\n标题：《致橡树》')
file.write('\n作者：舒婷')
file.write('\n时间：1977年3月')
file.close()
'''


'''
file = None
try:
    file = open('致橡树.txt', 'r', encoding='utf-8')
    print(file.read())
except FileNotFoundError:
    print('无法打开指定的文件!')
except LookupError:
    print('指定了未知的编码!')
except UnicodeDecodeError:
    print('读取文件时解码错误!')
finally:
    if file:
        file.close()
'''


'''
class InputError(ValueError):
    """自定义异常类型"""
    pass


def fac(num):
    """求阶乘"""
    if num < 0:
        raise InputError('只能计算非负整数的阶乘')
    if num in (0, 1):
        return 1
    return num * fac(num - 1)

flag = True
while flag:
    num = int(input('n = '))
    try:
        print(f'{num}! = {fac(num)}')
        flag = False
    except InputError as err:
        print(err)
'''


'''
try:
    with open('致橡树.txt', 'r', encoding='utf-8') as file:
        print(file.read())
except FileNotFoundError:
    print('无法打开指定的文件!')
except LookupError:
    print('指定了未知的编码!')
except UnicodeDecodeError:
    print('读取文件时解码错误!')
'''


'''
import json

my_dict = {
    'name': '骆昊',
    'age': 40,
    'friends': ['王大锤', '白元芳'],
    'cars': [
        {'brand': 'BMW', 'max_speed': 240},
        {'brand': 'Audi', 'max_speed': 280},
        {'brand': 'Benz', 'max_speed': 280}
    ]
}
with open('data.json', 'w') as file:
    json.dump(my_dict, file)

with open('data.json', 'r') as file:
    my_dict = json.load(file)
    print(type(my_dict))
    print(my_dict)
'''


'''
import requests

response = requests.get('https://apis.tianapi.com/guonei/index?key=99a9159476b792afa46184935a66047e&num=10')
if response.status_code == 200:
    data_model = response.json()
    for news in data_model['result']['newslist']:
        print(news['title'])
        print(news['url'])
        print('-' * 60)
'''


'''
import csv
import random

with open('scores.csv', 'w') as file:
    writer = csv.writer(file, delimiter='|', quoting=csv.QUOTE_ALL)
    writer.writerow(['姓名', '语文', '数学', '英语'])
    names = ['关羽', '张飞', '赵云', '马超', '黄忠']
    for name in names:
        scores = [random.randrange(50, 101) for _ in range(3)]
        scores.insert(0, name)
        writer.writerow(scores)
'''


'''
import csv

with open('scores.csv', 'r') as file:
    reader = csv.reader(file, delimiter='|')
    for data_list in reader:
        print(reader.line_num, end='\t')
        for elem in data_list:
            print(elem, end='\t')
        print()
'''


'''
import xlrd

# 使用xlrd模块的open_workbook函数打开指定Excel文件并获得Book对象（工作簿）
wb = xlrd.open_workbook('阿里巴巴2020年股票数据.xls')
# 通过Book对象的sheet_names方法可以获取所有表单名称
sheetnames = wb.sheet_names()
print(sheetnames)
# 通过指定的表单名称获取Sheet对象（工作表）
sheet = wb.sheet_by_name(sheetnames[0])
# 通过Sheet对象的nrows和ncols属性获取表单的行数和列数
print(sheet.nrows, sheet.ncols)
for row in range(sheet.nrows):
    for col in range(sheet.ncols):
        # 通过Sheet对象的cell方法获取指定Cell对象（单元格）
        # 通过Cell对象的value属性获取单元格中的值
        value = sheet.cell(row, col).value
        # 对除首行外的其他行进行数据格式化处理
        if row > 0:
            # 第1列的xldate类型先转成元组再格式化为“年月日”的格式
            if col == 0:
                # xldate_as_tuple函数的第二个参数只有0和1两个取值
                # 其中0代表以1900-01-01为基准的日期，1代表以1904-01-01为基准的日期
                value = xlrd.xldate_as_tuple(value, 0)
                value = f'{value[0]}年{value[1]:>02d}月{value[2]:>02d}日'
            # 其他列的number类型处理成小数点后保留两位有效数字的浮点数
            else:
                value = f'{value:.2f}'
        print(value, end='\t')
    print()
# 获取最后一个单元格的数据类型
# 0 - 空值，1 - 字符串，2 - 数字，3 - 日期，4 - 布尔，5 - 错误
last_cell_type = sheet.cell_type(sheet.nrows - 1, sheet.ncols - 1)
print(last_cell_type)
# 获取第一行的值（列表）
print(sheet.row_values(0))
# 获取指定行指定列范围的数据（列表）
# 第一个参数代表行索引，第二个和第三个参数代表列的开始（含）和结束（不含）索引
print(sheet.row_slice(3, 0, 5))
'''


'''
import random

import xlwt

student_names = ['关羽', '张飞', '赵云', '马超', '黄忠']
scores = [[random.randrange(50, 101) for _ in range(3)] for _ in range(5)]
# 创建工作簿对象（Workbook）
wb = xlwt.Workbook()
# 创建工作表对象（Worksheet）
sheet = wb.add_sheet('一年级二班')
# 添加表头数据
header_style = xlwt.XFStyle()
pattern = xlwt.Pattern()
pattern.pattern = xlwt.Pattern.SOLID_PATTERN
# 0 - 黑色、1 - 白色、2 - 红色、3 - 绿色、4 - 蓝色、5 - 黄色、6 - 粉色、7 - 青色
pattern.pattern_fore_colour = 5
header_style.pattern = pattern

font = xlwt.Font()
# 字体名称
font.name = '华文楷体'
# 字体大小（20是基准单位，18表示18px）
font.height = 20 * 18
# 是否使用粗体
font.bold = True
# 是否使用斜体
font.italic = False
# 字体颜色
font.colour_index = 1
header_style.font = font

titles = ('姓名', '语文', '数学', '英语')
for index, title in enumerate(titles):
    sheet.write(0, index, title, header_style)
# 将学生姓名和考试成绩写入单元格
for row in range(len(scores)):
    sheet.write(row + 1, 0, student_names[row])
    for col in range(len(scores[row])):
        sheet.write(row + 1, col + 1, scores[row][col])
# 保存Excel工作簿
wb.save('考试成绩表.xls')
'''


'''
import xlrd
import xlwt
from xlutils.copy import copy

# 打开文件用于读取
wb_for_read = xlrd.open_workbook('阿里巴巴2020年股票数据.xls')
sheet1 = wb_for_read.sheet_by_index(0)
nrows, ncols = sheet1.nrows, sheet1.ncols  # 获取行列数

# 复制工作簿用于写入
wb_for_write = copy(wb_for_read)
sheet2 = wb_for_write.get_sheet(0)

# 在最后一行下方添加公式
sheet2.write(nrows, 4, xlwt.Formula(f'average(E2:E{nrows})'))  # E列添加平均值
sheet2.write(nrows, 5, xlwt.Formula(f'sum(F2:F{nrows})'))      # F列添加总和

# 保存为新文件
wb_for_write.save('阿里巴巴2020年股票数据汇总.xls')
'''


'''
import datetime
import openpyxl

# 加载一个工作簿 ---> Workbook
wb = openpyxl.load_workbook('阿里巴巴2020年股票数据.xlsx')
# 获取工作表的名字
print(wb.sheetnames)
# 获取工作表 ---> Worksheet
sheet = wb.worksheets[0]
# 获得单元格的范围
print(sheet.dimensions)
# 获得行数和列数
print(sheet.max_row, sheet.max_column)

# 获取指定单元格的值
print(sheet.cell(3, 3).value)
print(sheet['C3'].value)
print(sheet['G255'].value)

# 获取多个单元格（嵌套元组）
print(sheet['A2:C5'])

# 读取所有单元格的数据
for row_ch in range(2, sheet.max_row + 1):
    for col_ch in 'ABCDEFG':
        value = sheet[f'{col_ch}{row_ch}'].value
        if type(value) == datetime.datetime:
            print(value.strftime('%Y年%m月%d日'), end='\t')
        elif type(value) == int:
            print(f'{value:<10d}', end='\t')
        elif type(value) == float:
            print(f'{value:.4f}', end='\t')
        else:
            print(value, end='\t')
    print()
'''


'''
import random

import openpyxl

# 第一步：创建工作簿（Workbook）
wb = openpyxl.Workbook()

# 第二步：添加工作表（Worksheet）
sheet = wb.active
sheet.title = '期末成绩'

titles = ('姓名', '语文', '数学', '英语')
for col_index, title in enumerate(titles):
    sheet.cell(1, col_index + 1, title)

names = ('关羽', '张飞', '赵云', '马超', '黄忠')
for row_index, name in enumerate(names):
    sheet.cell(row_index + 2, 1, name)
    for col_index in range(2, 5):
        sheet.cell(row_index + 2, col_index, random.randrange(50, 101))

# 第四步：保存工作簿
wb.save('考试成绩表.xlsx')
'''


'''
import openpyxl
from openpyxl.styles import Font, Alignment, Border, Side

# 对齐方式
alignment = Alignment(horizontal='center', vertical='center')
# 边框线条
side = Side(color='ff7f50', style='mediumDashed')

wb = openpyxl.load_workbook('考试成绩表.xlsx')
sheet = wb.worksheets[0]

# 调整行高和列宽
sheet.row_dimensions[1].height = 30
sheet.column_dimensions['E'].width = 120

sheet['E1'] = '平均分'
# 设置字体
sheet.cell(1, 5).font = Font(size=18, bold=True, color='ff1493', name='华文楷体')
# 设置对齐方式
sheet.cell(1, 5).alignment = alignment
# 设置单元格边框
sheet.cell(1, 5).border = Border(left=side, top=side, right=side, bottom=side)
for i in range(2, 7):
    # 公式计算每个学生的平均分
    sheet[f'E{i}'] = f'=average(B{i}:D{i})'
    sheet.cell(i, 5).font = Font(size=12, color='4169e1', italic=True)
    sheet.cell(i, 5).alignment = alignment

wb.save('考试成绩表.xlsx')
'''


'''
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference

wb = Workbook(write_only=True)
sheet = wb.create_sheet()

rows = [
    ('类别', '销售A组', '销售B组'),
    ('手机', 40, 30),
    ('平板', 50, 60),
    ('笔记本', 80, 70),
    ('外围设备', 20, 10),
]

# 向表单中添加行
for row in rows:
    sheet.append(row)

# 创建图表对象
chart = BarChart()
chart.type = 'col'
chart.style = 10
# 设置图表的标题
chart.title = '销售统计图'
# 设置图表纵轴的标题
chart.y_axis.title = '销量'
# 设置图表横轴的标题
chart.x_axis.title = '商品类别'
# 设置数据的范围
data = Reference(sheet, min_col=2, min_row=1, max_row=5, max_col=3)
# 设置分类的范围
cats = Reference(sheet, min_col=1, min_row=2, max_row=5)
# 给图表添加数据
chart.add_data(data, titles_from_data=True)
# 给图表设置分类
chart.set_categories(cats)
chart.shape = 4
# 将图表添加到表单指定的单元格中
sheet.add_chart(chart, 'A10')

wb.save('demo.xlsx')
'''


'''
from docx import Document
from docx.shared import Cm, Pt

from docx.document import Document as Doc

# 创建代表Word文档的Doc对象
document = Document()  # type: Doc
# 添加大标题
document.add_heading('快快乐乐学Python', 0)
# 添加段落
p = document.add_paragraph('Python是一门非常流行的编程语言，它')
run = p.add_run('简单')
run.bold = True
run.font.size = Pt(18)
p.add_run('而且')
run = p.add_run('优雅')
run.font.size = Pt(18)
run.underline = True
p.add_run('。')

# 添加一级标题
document.add_heading('Heading, level 1', level=1)
# 添加带样式的段落
document.add_paragraph('Intense quote', style='Intense Quote')
# 添加无序列表
document.add_paragraph(
    'first item in unordered list', style='List Bullet'
)
document.add_paragraph(
    'second item in ordered list', style='List Bullet'
)
# 添加有序列表
document.add_paragraph(
    'first item in ordered list', style='List Number'
)
document.add_paragraph(
    'second item in ordered list', style='List Number'
)

# 添加图片（注意路径和图片必须要存在）
document.add_picture('resources/guido.jpg', width=Cm(5.2))

# 添加分节符
document.add_section()

records = (
    ('骆昊', '男', '1995-5-5'),
    ('孙美丽', '女', '1992-2-2')
)
# 添加表格
table = document.add_table(rows=1, cols=3)
table.style = 'Dark List'
hdr_cells = table.rows[0].cells
hdr_cells[0].text = '姓名'
hdr_cells[1].text = '性别'
hdr_cells[2].text = '出生日期'
# 为表格添加行
for name, sex, birthday in records:
    row_cells = table.add_row().cells
    row_cells[0].text = name
    row_cells[1].text = sex
    row_cells[2].text = birthday

# 添加分页符
document.add_page_break()

# 保存文档
document.save('demo.docx')
'''


'''
from docx import Document
from docx.document import Document as Doc

# 将真实信息用字典的方式保存在列表中
employees = [
    {
        'name': '骆昊',
        'id': '100200198011280001',
        'sdate': '2008年3月1日',
        'edate': '2012年2月29日',
        'department': '产品研发',
        'position': '架构师',
        'company': '成都华为技术有限公司'
    },
    {
        'name': '王大锤',
        'id': '510210199012125566',
        'sdate': '2019年1月1日',
        'edate': '2021年4月30日',
        'department': '产品研发',
        'position': 'Python开发工程师',
        'company': '成都谷道科技有限公司'
    },
    {
        'name': '李元芳',
        'id': '2102101995103221599',
        'sdate': '2020年5月10日',
        'edate': '2021年3月5日',
        'department': '产品研发',
        'position': 'Java开发工程师',
        'company': '同城企业管理集团有限公司'
    },
]
# 对列表进行循环遍历，批量生成Word文档
for emp_dict in employees:
    # 读取离职证明模板文件
    doc = Document('离职证明模板.docx')  # type: Doc
    # 循环遍历所有段落寻找占位符
    for p in doc.paragraphs:
        if '{' not in p.text:
            continue
        # 不能直接修改段落内容，否则会丢失样式
        # 所以需要对段落中的元素进行遍历并进行查找替换
        for run in p.runs:
            if '{' not in run.text:
                continue
            # 将占位符换成实际内容
            start, end = run.text.find('{'), run.text.find('}')
            key, place_holder = run.text[start + 1:end], run.text[start:end + 1]
            run.text = run.text.replace(place_holder, emp_dict[key])
    # 每个人对应保存一个Word文档
    doc.save(f'{emp_dict["name"]}离职证明.docx')
'''


'''
from pptx import Presentation

# 创建幻灯片对象
pres = Presentation()

# 选择母版添加一页
title_slide_layout = pres.slide_layouts[0]
slide = pres.slides.add_slide(title_slide_layout)
# 获取标题栏和副标题栏
title = slide.shapes.title
subtitle = slide.placeholders[1]
# 编辑标题和副标题
title.text = "Welcome to Python"
subtitle.text = "Life is short, I use Python"

# 选择母版添加一页
bullet_slide_layout = pres.slide_layouts[1]
slide = pres.slides.add_slide(bullet_slide_layout)
# 获取页面上所有形状
shapes = slide.shapes
# 获取标题和主体
title_shape = shapes.title
body_shape = shapes.placeholders[1]
# 编辑标题
title_shape.text = 'Introduction'
# 编辑主体内容
tf = body_shape.text_frame
tf.text = 'History of Python'
# 添加一个一级段落
p = tf.add_paragraph()
p.text = 'X\'max 1989'
p.level = 1
# 添加一个二级段落
p = tf.add_paragraph()
p.text = 'Guido began to write interpreter for Python.'
p.level = 2

# 保存幻灯片
pres.save('test.pptx')




import re

sentence = 'Oh, shit! 你是傻逼吗? Fuck you.'
purified = re.sub('fuck|shit|[傻煞沙][比笔逼叉缺吊碉雕]',
                  '*', sentence, flags=re.IGNORECASE)
print(purified)  # Oh, *! 你是*吗? * you.




import re

poem = '窗前明月光，疑是地上霜。举头望明月，低头思故乡。'
sentences_list = re.split(r'[，。]', poem)
sentences_list = [sentence for sentence in sentences_list if sentence]
for sentence in sentences_list:
    print(sentence)
'''



